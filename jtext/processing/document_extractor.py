"""
Document text extraction for various formats.

This module provides text extraction capabilities for PDF, DOCX, PPTX,
and HTML documents with metadata preservation and quality metrics.
"""

import time
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
import psutil

# Document processing imports
import fitz  # PyMuPDF for PDF
from docx import Document
from pptx import Presentation
import html2text

from ..utils.logging import get_logger
from ..utils.validation import validate_document_file
from ..core.ocr_hybrid import ProcessingResult


logger = get_logger(__name__)


class DocumentExtractor:
    """
    Multi-format document text extraction.

    This class provides text extraction from various document formats
    including PDF, DOCX, PPTX, and HTML with metadata preservation.
    """

    def __init__(self):
        """Initialize the document extractor."""
        logger.debug("Initialized DocumentExtractor")

    def extract_text(self, file_path: str) -> ProcessingResult:
        """
        Extract text from a document file.

        Args:
            file_path: Path to the document file

        Returns:
            ProcessingResult containing extracted text and metadata

        Raises:
            ValueError: If file format is not supported
            FileNotFoundError: If file doesn't exist
        """
        start_time = time.time()
        start_memory = psutil.Process().memory_info().rss / 1024 / 1024  # MB

        logger.info(f"Starting document extraction for: {file_path}")

        # Validate input file
        if not validate_document_file(file_path):
            raise ValueError(f"Invalid document file: {file_path}")

        try:
            file_path_obj = Path(file_path)
            file_extension = file_path_obj.suffix.lower()

            # Route to appropriate extractor based on file type
            if file_extension == ".pdf":
                text, metadata = self._extract_pdf(file_path)
            elif file_extension in [".docx"]:
                text, metadata = self._extract_docx(file_path)
            elif file_extension in [".pptx"]:
                text, metadata = self._extract_pptx(file_path)
            elif file_extension in [".html", ".htm"]:
                text, metadata = self._extract_html(file_path)
            else:
                raise ValueError(f"Unsupported document format: {file_extension}")

            # Calculate processing metrics
            processing_time = time.time() - start_time
            end_memory = psutil.Process().memory_info().rss / 1024 / 1024  # MB
            memory_usage = end_memory - start_memory

            # Create result object
            result = ProcessingResult(
                source_path=file_path,
                text=text,
                confidence=1.0,  # Document extraction is deterministic
                processing_time=processing_time,
                memory_usage=memory_usage,
                corrections_applied=0,
                correction_ratio=0.0,
            )

            # Add document-specific metadata
            result.document_metadata = metadata

            logger.info(
                f"Document extraction completed in {processing_time:.2f}s, "
                f"extracted {len(text)} characters"
            )

            return result

        except Exception as e:
            logger.error(f"Document extraction failed for {file_path}: {e}")
            raise

    def _extract_pdf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from PDF document.

        Args:
            file_path: Path to PDF file

        Returns:
            Tuple of (extracted_text, metadata)
        """
        logger.debug("Extracting text from PDF")

        text_parts = []
        metadata = {
            "format": "pdf",
            "pages": 0,
            "has_images": False,
            "has_tables": False,
        }

        try:
            doc = fitz.open(file_path)
            metadata["pages"] = len(doc)

            # Add document metadata if available
            if hasattr(doc, "metadata") and doc.metadata:
                for key, value in doc.metadata.items():
                    if value:  # Only add non-empty values
                        metadata[key] = value

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                text_parts.append(page_text)

                # Check for images and tables
                if page.get_images():
                    metadata["has_images"] = True
                if page.get_text("dict").get("blocks"):
                    metadata["has_tables"] = True

            doc.close()

            full_text = "\n\n".join(text_parts)
            return full_text, metadata

        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            raise

    def _extract_docx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from DOCX document.

        Args:
            file_path: Path to DOCX file

        Returns:
            Tuple of (extracted_text, metadata)
        """
        logger.debug("Extracting text from DOCX")

        metadata = {
            "format": "docx",
            "paragraphs": 0,
            "tables": 0,
            "has_images": False,
        }

        try:
            doc = Document(file_path)

            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
                    metadata["paragraphs"] += 1

            # Extract table content
            for table in doc.tables:
                metadata["tables"] += 1
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))

            # Check for images
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    metadata["has_images"] = True
                    break

            full_text = "\n\n".join(text_parts)
            return full_text, metadata

        except Exception as e:
            logger.error(f"DOCX extraction failed: {e}")
            raise

    def _extract_pptx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from PPTX presentation.

        Args:
            file_path: Path to PPTX file

        Returns:
            Tuple of (extracted_text, metadata)
        """
        logger.debug("Extracting text from PPTX")

        metadata = {
            "format": "pptx",
            "slides": 0,
            "has_images": False,
            "has_tables": False,
        }

        try:
            prs = Presentation(file_path)
            metadata["slides"] = len(prs.slides)

            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    elif hasattr(shape, "table"):
                        metadata["has_tables"] = True
                        # Extract table content
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                    elif hasattr(shape, "image"):
                        metadata["has_images"] = True

                if len(slide_text) > 1:  # More than just the slide header
                    text_parts.append("\n".join(slide_text))

            full_text = "\n\n".join(text_parts)
            return full_text, metadata

        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise

    def _extract_html(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from HTML document.

        Args:
            file_path: Path to HTML file

        Returns:
            Tuple of (extracted_text, metadata)
        """
        logger.debug("Extracting text from HTML")

        metadata = {
            "format": "html",
            "has_images": False,
            "has_links": False,
        }

        try:
            with open(file_path, "r", encoding="utf-8") as f:
                html_content = f.read()

            # Check for images and links
            if "<img" in html_content.lower():
                metadata["has_images"] = True
            if "<a href" in html_content.lower():
                metadata["has_links"] = True

            # Convert HTML to text
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = False
            h.body_width = 0  # Don't wrap lines

            text = h.handle(html_content)
            return text, metadata

        except Exception as e:
            logger.error(f"HTML extraction failed: {e}")
            raise

    def extract_multiple(self, file_paths: List[str]) -> List[ProcessingResult]:
        """
        Extract text from multiple document files.

        Args:
            file_paths: List of document file paths

        Returns:
            List of ProcessingResult objects
        """
        results = []

        for file_path in file_paths:
            try:
                result = self.extract_text(file_path)
                results.append(result)
                logger.info(f"Successfully extracted: {file_path}")
            except Exception as e:
                logger.error(f"Failed to extract {file_path}: {e}")
                # Continue with other files
                continue

        logger.info(f"Extracted {len(results)}/{len(file_paths)} documents")
        return results
